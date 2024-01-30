from fastapi import FastAPI
# ... 其他导入 ...
from fastapi import UploadFile, File, HTTPException
import os
from fastapi.responses import JSONResponse
from fastapi import Query
from typing import List
import uvicorn
from io import BytesIO
import concurrent.futures

from transformers import BertTokenizer, BertModel
import torch

import pickle  # 添加导入语句

from PyPDF2 import PdfReader

# 引入您原有的逻辑
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import docx
import pptx
import markdown
import json

import mimetypes


import sqlite3
from typing import List, Tuple
from collections import defaultdict

# 全局变量来存储缓存的文档和索引
docs_cache = defaultdict(dict)  # 使用 defaultdict 来自动为新的 ID 创建字典

# 数据库名
db_file = 'files.db'

# 数据库初始化，需要另外创建一个文件执行去创建数据库
def init_db():
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY,
            filename TEXT,
            file_type TEXT,
            content BLOB,
            vector BLOB
            text_keys BLOB
        )
    ''')
    conn.commit()
    conn.close()

# 保存文件、向量和文本键列表到数据库
def save_file_to_db(filename, file_type, content, vector, text_keys):
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO files (filename, file_type, content, vector, text_keys)
        VALUES (?, ?, ?, ?, ?)
    ''', (filename, file_type, content, vector, text_keys))
    conn.commit()
    conn.close()

# 从数据库获取所有文件的向量
def get_all_vectors():
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute('SELECT content, vector FROM files')
    rows = cursor.fetchall()
    conn.close()
    # 确保正确解码向量
    return [(content, np.frombuffer(vector, dtype='float32')) for content, vector in rows]




# 处理函数
# 使用Sentence Transformers获取文本嵌入
class SentenceEmbeddings:
    def __init__(self, model_name='sentence-transformers/paraphrase-multilingual-mpnet-base-v2'):
        # 使用提供的模型名称或路径，如果没有提供则使用默认值
        self.model = SentenceTransformer(model_name)

    def get_embedding(self, text):
        return self.model.encode(text, show_progress_bar=False)





# 使用BERT获取文本嵌入
class BertEmbeddings:
    def __init__(self):
        self.tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')
        self.model = BertModel.from_pretrained('bert-base-uncased')

    def get_embedding(self, text):
        inputs = self.tokenizer(text, return_tensors='pt', truncation=True, max_length=512)
        with torch.no_grad():
            outputs = self.model(**inputs)
        return np.squeeze(outputs.last_hidden_state.mean(dim=1).numpy())
    
# 
def get_pdf_text(pdf_bytes):
    reader = PdfReader(BytesIO(pdf_bytes))
    text = []
    for page in reader.pages:
        text.append(page.extract_text())
    return text

# 获取文档文本，支持多种格式
def get_document_text(docs, file_type):
    text_blocks = []
    for doc in docs:
        if file_type == "application/pdf":
            pdf_reader = PdfReader(BytesIO(doc))
            for page in pdf_reader.pages:
                text_blocks.extend(page.extract_text().split("\n\n"))
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            docx_file = docx.Document(BytesIO(doc))
            text_blocks.extend([para.text for para in docx_file.paragraphs if para.text])
        elif file_type == "text/plain":
            text_blocks.append(doc.decode("utf-8"))
        elif file_type == "application/json":
            json_data = json.loads(doc.decode("utf-8"))
            # 递归地抽取所有字符串值
            def extract_text_from_json(obj):
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        extract_text_from_json(v)
                elif isinstance(obj, list):
                    for item in obj:
                        extract_text_from_json(item)
                elif isinstance(obj, str):
                    text_blocks.append(obj)
            extract_text_from_json(json_data)
        elif file_type == "text/markdown":
            md_text = markdown.markdown(doc.decode("utf-8"))
            text_blocks.extend(md_text.split("\n\n"))
        elif file_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            ppt = pptx.Presentation(BytesIO(doc))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_blocks.append(shape.text)
    return text_blocks

def process_file_from_db(file_blob, file_type):
    if file_type == "application/pdf":
        # 对于 PDF 文件，使用 PyPDF2 提取文本
        pdf_reader = PdfReader(BytesIO(file_blob))
        text_blocks = [page.extract_text() for page in pdf_reader.pages if page.extract_text()]
        return text_blocks
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        # 对于 Word 文件，使用 docx 提取文本
        docx_file = docx.Document(BytesIO(file_blob))
        text_blocks = [para.text for para in docx_file.paragraphs if para.text]
        return text_blocks
    elif file_type == "text/plain":
        # 对于普通文本文件，直接读取
        text_blocks = [file_blob.decode("utf-8")]
        return text_blocks
    elif file_type == "application/json":
        # 对于 JSON 文件，使用 json 提取文本
        json_data = json.loads(file_blob.decode("utf-8"))
        # 递归地抽取所有字符串值
        def extract_text_from_json(obj):
            if isinstance(obj, dict):
                for k, v in obj.items():
                    extract_text_from_json(v)
            elif isinstance(obj, list):
                for item in obj:
                    extract_text_from_json(item)
            elif isinstance(obj, str):
                text_blocks.append(obj)
        extract_text_from_json(json_data)
        return text_blocks
    elif file_type == "text/markdown":
        # 对于 Markdown 文件，使用 markdown 提取文本
        md_text = markdown.markdown(file_blob.decode("utf-8"))
        text_blocks = md_text.split("\n\n")
        return text_blocks
    elif file_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            ppt = pptx.Presentation(BytesIO(doc))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_blocks.append(shape.text)
    return text_blocks
    

# 将文本转换为嵌入向量字典
def text_to_vector_dict(text_blocks,model_name,search_mode="cosine"):
    if search_mode == "cosine":
        sentence_embeddings = SentenceEmbeddings(model_name)
        text_to_vector = {}
        with concurrent.futures.ThreadPoolExecutor() as executor:
            embeddings = list(executor.map(sentence_embeddings.get_embedding, text_blocks))
        # 归一化嵌入
        embeddings_normalized = [embedding / np.linalg.norm(embedding) for embedding in embeddings if np.linalg.norm(embedding) > 0]
        text_to_vector = dict(zip(text_blocks, embeddings_normalized))
    elif search_mode == "euclidean":
        sentence_embeddings = SentenceEmbeddings(model_name)
        text_to_vector = {}
        with concurrent.futures.ThreadPoolExecutor() as executor:
            embeddings = list(executor.map(sentence_embeddings.get_embedding, text_blocks))
        text_to_vector = dict(zip(text_blocks, embeddings))
    else:
        print("Invalid search mode. Please choose 'cosine' or 'euclidean'.")
        return
    return text_to_vector

# 将文本转换为嵌入向量字典 余弦相似度
# def text_to_vector_dict(text_blocks,model_name):
#     sentence_embeddings = SentenceEmbeddings(model_name)
#     text_to_vector = {}
#     with concurrent.futures.ThreadPoolExecutor() as executor:
#         embeddings = list(executor.map(sentence_embeddings.get_embedding, text_blocks))
#     # 归一化嵌入
#     embeddings_normalized = [embedding / np.linalg.norm(embedding) for embedding in embeddings if np.linalg.norm(embedding) > 0]
#     text_to_vector = dict(zip(text_blocks, embeddings_normalized))
#     return text_to_vector


# 创建FAISS索引
def create_faiss_index(vector_dict):
    dimension = len(next(iter(vector_dict.values())))
    index = faiss.IndexFlatL2(dimension)
    vectors = np.array(list(vector_dict.values())).astype('float32')
    index.add(vectors)
    return index, list(vector_dict.keys())

# 在FAISS索引中搜索
def search_faiss_index(index, query_vector, text_keys, k=10,search_mode="cosine"):
    if search_mode == "cosine":
        # 确保查询向量归一化
        norm_query_vector = query_vector / np.linalg.norm(query_vector)
        query_vector = np.array([norm_query_vector]).astype('float32')
        distances, indices = index.search(query_vector, k)
        results = []
        for i in range(len(indices[0])):
            idx = indices[0][i]
            if idx < len(text_keys):
                results.append((text_keys[idx], distances[0][i]))
    elif search_mode == "euclidean":
        query_vector = np.array([query_vector]).astype('float32')
        distances, indices = index.search(query_vector, k)
        results = []
        for i in range(len(indices[0])):  # 确保不超过索引界限
            idx = indices[0][i]
            if idx < len(text_keys):  # 检查索引是否有效
                results.append((text_keys[idx], distances[0][i]))
    else:
        print("Invalid search mode. Please choose 'cosine' or 'euclidean'.")
        return
    return results

# 在FAISS索引中搜索 余弦相似度
# def search_faiss_index(index, query_vector, text_keys, k=5):
#     # 确保查询向量归一化
#     norm_query_vector = query_vector / np.linalg.norm(query_vector)
#     query_vector = np.array([norm_query_vector]).astype('float32')
#     distances, indices = index.search(query_vector, k)
#     results = []
#     for i in range(len(indices[0])):
#         idx = indices[0][i]
#         if idx < len(text_keys):
#             results.append((text_keys[idx], distances[0][i]))
#     return results

def get_all_files_from_db() -> List[Tuple[bytes, str]]:
    """
    从数据库中检索所有文件的内容和类型。

    返回:
        List[Tuple[bytes, str]]: 返回一个元组列表，每个元组包含文件的二进制内容和文件类型。
    """
    conn = sqlite3.connect('files.db')
    cursor = conn.cursor()
    try:
        # 选择 content 和 file_type 字段
        cursor.execute('SELECT content, file_type FROM files')
        rows = cursor.fetchall()
        return rows
    except sqlite3.Error as e:
        print(f"数据库错误: {e}")
        return []
    finally:
        conn.close()


def get_files_from_db_by_ids(ids: List[int]) -> List[Tuple[bytes, str]]:
    conn = sqlite3.connect('files.db')
    cursor = conn.cursor()
    try:
        # 构建查询字符串
        placeholders = ', '.join('?' for _ in ids)
        query = f'SELECT content, file_type FROM files WHERE id IN ({placeholders})'
        cursor.execute(query, ids)
        rows = cursor.fetchall()
        return rows
    except sqlite3.Error as e:
        print(f"数据库错误: {e}")
        return []
    finally:
        conn.close()

# 格式化文本以用于GPT
def format_text_for_gpt(text):
    clean_text = text.replace('\n', ' ').strip()
    return clean_text


# 加载 FAISS 索引和文本键列表
# 加载 FAISS 索引和文本键列表
def load_faiss_index(file_name):
    conn = sqlite3.connect('files.db')
    cursor = conn.cursor()
    cursor.execute('SELECT vector, text_keys FROM files WHERE filename=?', (file_name,))
    row = cursor.fetchone()
    conn.close()

    if row:
        vector_blob = row[0]
        text_keys = json.loads(row[1])

        # 使用 pickle 反序列化向量
        vector = pickle.loads(vector_blob)

        faiss_index = faiss.deserialize_index(vector)
        return faiss_index, text_keys
    else:
        return None, None

    

# 全局变量来存储文档和索引
docs_text_keys = None
faiss_index = None


# 添加全局变量来存储文档和索引
upload_dir = "uploaded_files"  # 这里设置为默认的上传目录

# docfunc_router = FastAPI()

def include_docfunc_router(app: FastAPI):
    # 修改 /upload-files/ 的处理函数，将 upload_dir 作为接口参数传递
    @app.post("/upload-files/", tags=["文档向量检索"])
    async def upload_files(
        files: List[UploadFile] = File(...),
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        search_mode: str = Query("cosine"),
        upload_dir: str = Query("uploaded_files")  # 将 upload_dir 作为接口参数
    ):
        global docs_text_keys, faiss_index

        try:
            # 遍历所有文件
            for file in files:
                # 创建以文件名命名的文件夹
                file_folder = os.path.join(upload_dir, os.path.splitext(file.filename)[0])
                os.makedirs(file_folder, exist_ok=True)

                # 存储所有文档的文本
                all_docs_text = []

                # 提取文本
                contents = await file.read()
                file_type = file.content_type
                text = get_document_text([contents], file_type)  # 直接传递 contents 和 file_type
                all_docs_text.extend(text)

                # 将文件内容写入本地文件
                file_path = os.path.join(file_folder, file.filename)
                with open(file_path, "wb") as local_file:
                    local_file.write(contents)

                # 将文本转换为嵌入向量并创建 FAISS 索引
                vector_dict = text_to_vector_dict(all_docs_text, model_name, search_mode)
                faiss_index, docs_text_keys = create_faiss_index(vector_dict)

                # 保存当前构建索引的文本键列表到文件夹，文件名与源文件名一样
                keys_file_name = f"{os.path.splitext(file.filename)[0]}_text_keys.json"
                keys_file_path = os.path.join(file_folder, keys_file_name)
                with open(keys_file_path, "w") as keys_file:
                    json.dump(docs_text_keys, keys_file)

                # 保存 FAISS 索引到磁盘，使用上传文件的名称作为索引文件的名称
                index_file_name = f"{os.path.splitext(file.filename)[0]}_faiss_index.index"
                index_file_path = os.path.join(file_folder, index_file_name)
                faiss.write_index(faiss_index, index_file_path)

            return JSONResponse(content={"status": "Processing completed successfully", "uploaded_files": [file.filename for file in files]})
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)



    @app.post("/search/", tags=["文档向量检索"])
    async def search_docs(
        query: str,
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        k: int = Query(5),
        search_mode: str = Query("cosine"),
        search_dir: str = Query("uploaded_files"),
        file_names: List[str] = Query(["BS网管架构技术介绍V1.0.pdf", "其他文件名.pdf"])  # 修改为文件名列表
    ):
        global docs_text_keys, faiss_index

        try:
            results_list = []

            for file_name in file_names:
                # 构建文件夹路径
                file_folder = os.path.join(search_dir, os.path.splitext(file_name)[0])

                if not faiss_index or not docs_text_keys:
                    # 加载本地保存的FAISS索引
                    index_file_name = f"{os.path.splitext(file_name)[0]}_faiss_index.index"
                    index_file_path = os.path.join(file_folder, index_file_name)
                    faiss_index = faiss.read_index(index_file_path)

                    # 加载本地保存的文本键列表
                    keys_file_name = f"{os.path.splitext(file_name)[0]}_text_keys.json"
                    keys_file_path = os.path.join(file_folder, keys_file_name)
                    with open(keys_file_path, "r") as keys_file:
                        docs_text_keys = json.load(keys_file)

                    if not faiss_index or not docs_text_keys:
                        raise HTTPException(status_code=500, detail="Failed to load FAISS index or text keys: Index or keys are None")

                # 初始化 SentenceEmbeddings 实例
                sentence_embeddings = SentenceEmbeddings(model_name)

                # 将查询转换为向量
                query_vector = sentence_embeddings.get_embedding(query)

                # 在 FAISS 索引中进行搜索
                results = search_faiss_index(faiss_index, query_vector, docs_text_keys, k, search_mode)

                # 格式化搜索结果
                formatted_results = [{"text": text, "distance": float(distance)} for text, distance in results]

                 # 将所有"text"的值拼接在一起
                combined_text = '\n'.join(result['text'] for result in formatted_results)
                # return {"combined_text": combined_text}

                # 清除检索缓存
                docs_text_keys = None
                faiss_index = None

                results_list.append({file_name: combined_text})

            return {"results": results_list}
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)


    @app.post("/upload-files-db/", tags=["文档向量检索"])
    async def upload_files_db(
        files: List[UploadFile] = File(...),
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        search_mode: str = Query("cosine")
    ):
        global docs_text_keys, faiss_index

        try:
            # 遍历所有文件
            for file in files:
                contents = await file.read()
                file_type = file.content_type
                text = get_document_text([contents], file_type)  # 直接传递 contents 和 file_type

                # 将文本转换为嵌入向量并创建 FAISS 索引
                vector_dict = text_to_vector_dict(text, model_name, search_mode)
                faiss_index, docs_text_keys = create_faiss_index(vector_dict)

                # 保存向量索引和文本键列表到数据库
                index_file = faiss.serialize_index(faiss_index)
                text_keys_json = json.dumps(docs_text_keys)
                save_file_to_db(file.filename, "faiss", contents, index_file, text_keys_json)

            return JSONResponse(content={"status": "Processing completed successfully", "uploaded_files": [file.filename for file in files]})
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)

    @app.post("/search-db/", tags=["文档向量检索"])
    async def search_docs_db(
        query: str,
        ids: List[int] = Query(None),
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        k: int = Query(5),
        search_mode: str = Query("cosine")
    ):
        try:
            print("Requested IDs:", ids)
            # 检查是否提供了 IDs
            if ids is None:
                return JSONResponse(content={"status": "No IDs provided"}, status_code=400)

            results_list = []

            # 遍历提供的 IDs
            for doc_id in ids:
                # 加载数据库中的索引和文本键列表
                conn = sqlite3.connect('files.db')
                cursor = conn.cursor()
                cursor.execute('SELECT filename FROM files WHERE id=?', (doc_id,))
                row = cursor.fetchone()
                conn.close()

                if row:
                    file_name = row[0]
                    faiss_index, docs_text_keys = load_faiss_index(file_name)

                    if faiss_index and docs_text_keys:
                        # 初始化 SentenceEmbeddings 实例
                        sentence_embeddings = SentenceEmbeddings(model_name)

                        # 将查询转换为向量
                        query_vector = sentence_embeddings.get_embedding(query)

                        # 在 FAISS 索引中进行搜索
                        results = search_faiss_index(faiss_index, query_vector, docs_text_keys, k, search_mode)

                        # 格式化搜索结果
                        formatted_results = [{"text": text, "distance": float(distance)} for text, distance in results]

                        # 将所有 "text" 的值拼接在一起
                        combined_text = '\n'.join(result['text'] for result in formatted_results)

                        results_list.append({file_name: combined_text})
                    else:
                        return JSONResponse(content={"status": f"Failed to load index for file with ID {doc_id}"}, status_code=500)

            return {"results": results_list}
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)

    @app.post("/upload-file-from-path-db/", tags=["文档向量检索"])
    async def upload_files_from_path_db(file_paths: List[str] = Query(...), model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"), search_mode: str = Query("cosine")):
        # 存储所有文档的文本
        all_docs_text = []

        try:
            for file_path in file_paths:
                try:
                    with open(file_path, "rb") as file:
                        contents = file.read()
                        file_type = "application/pdf" if file_path.endswith('.pdf') else \
                                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document" if file_path.endswith('.docx') else \
                                    "text/plain" if file_path.endswith('.txt') else \
                                    "application/json" if file_path.endswith('.json') else \
                                    "text/markdown" if file_path.endswith('.md') else \
                                    "application/vnd.ms-powerpoint" if file_path.endswith('.ppt') else \
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation" if file_path.endswith('.pptx') else \
                                    "unknown"
                        text = get_document_text([contents], file_type)
                        all_docs_text.extend(text)

                        # 将文本转换为向量并存储在数据库中
                        vector_dict = text_to_vector_dict(text, model_name, search_mode)
                        vector = np.mean(list(vector_dict.values()), axis=0)
                        # 使用文件路径中的文件名作为数据库中的文件名
                        save_file_to_db(file_path.split("/")[-1], file_type, contents, vector.tobytes())
                except Exception as file_error:
                    raise HTTPException(status_code=500, detail=f"Failed to read file {file_path}: {str(file_error)}")

            return JSONResponse(content={"status": "File uploaded and processed successfully"})
        except Exception as e:
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)

    @app.post("/upload-file-from-path/", tags=["文档向量检索"])
    async def upload_files_from_path(
        file_paths: List[str] = Query(...),
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        search_mode: str = Query("cosine")
    ):
        global docs_text_keys, faiss_index

        try:
            # 遍历所有文件路径
            for file_path in file_paths:
                # 创建以文件名命名的文件夹
                file_folder = os.path.join(os.path.dirname(file_path), os.path.splitext(os.path.basename(file_path))[0])
                os.makedirs(file_folder, exist_ok=True)

                # 存储所有文档的文本
                all_docs_text = []

                # 读取文件内容
                with open(file_path, "rb") as file:
                    contents = file.read()

                # 提取文本
                file_type = mimetypes.guess_type(file_path)[0]
                text = get_document_text([contents], file_type)  # 直接传递 contents 和 file_type
                all_docs_text.extend(text)

                # 将文件内容写入本地文件
                file_name = os.path.basename(file_path)
                # local_file_path = os.path.join(file_folder, file_name)
                # with open(local_file_path, "wb") as local_file:
                #     local_file.write(contents)

                # 将文本转换为嵌入向量并创建 FAISS 索引
                vector_dict = text_to_vector_dict(all_docs_text, model_name, search_mode)
                faiss_index, docs_text_keys = create_faiss_index(vector_dict)

                # 保存当前构建索引的文本键列表到文件夹，文件名与源文件名一样
                keys_file_name = f"{os.path.splitext(file_name)[0]}_text_keys.json"
                keys_file_path = os.path.join(file_folder, keys_file_name)
                with open(keys_file_path, "w") as keys_file:
                    json.dump(docs_text_keys, keys_file)

                # 保存 FAISS 索引到磁盘，使用上传文件的名称作为索引文件的名称
                index_file_name = f"{os.path.splitext(file_name)[0]}_faiss_index.index"
                index_file_path = os.path.join(file_folder, index_file_name)
                faiss.write_index(faiss_index, index_file_path)

            return JSONResponse(content={"status": "Processing completed successfully", "uploaded_files": file_paths})
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)
        
    @app.post("/search-path/", tags=["文档向量检索"])
    async def search_docs_from_path(
        query: str,
        model_name: str = Query("sentence-transformers/paraphrase-multilingual-mpnet-base-v2"),
        k: int = Query(5),
        search_mode: str = Query("cosine"),
        file_paths: List[str] = Query(...),  # 修改为文件路径列表
    ):
        global docs_text_keys, faiss_index

        try:
            results_list = []

            for file_path in file_paths:
                # 构建文件夹路径，以文件名命名
                file_folder = os.path.join(os.path.dirname(file_path), os.path.splitext(os.path.basename(file_path))[0])

                if not faiss_index or not docs_text_keys:
                    # 加载本地保存的FAISS索引
                    index_file_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_faiss_index.index"
                    index_file_path = os.path.join(file_folder, index_file_name)
                    faiss_index = faiss.read_index(index_file_path)

                    # 加载本地保存的文本键列表
                    keys_file_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_text_keys.json"
                    keys_file_path = os.path.join(file_folder, keys_file_name)
                    with open(keys_file_path, "r") as keys_file:
                        docs_text_keys = json.load(keys_file)

                    if not faiss_index or not docs_text_keys:
                        raise HTTPException(status_code=500, detail="Failed to load FAISS index or text keys: Index or keys are None")

                # 初始化 SentenceEmbeddings 实例
                sentence_embeddings = SentenceEmbeddings(model_name)

                # 将查询转换为向量
                query_vector = sentence_embeddings.get_embedding(query)

                # 在 FAISS 索引中进行搜索
                results = search_faiss_index(faiss_index, query_vector, docs_text_keys, k, search_mode)

                # 格式化搜索结果
                formatted_results = [{"text": text, "distance": float(distance)} for text, distance in results]

                # 将所有"text"的值拼接在一起
                combined_text = '\n'.join(result['text'] for result in formatted_results)

                # 清除检索缓存
                docs_text_keys = None
                faiss_index = None

                results_list.append({os.path.basename(file_path): combined_text})

            return {"results": results_list}
        except Exception as e:
            # 处理失败时返回带有错误消息的 JSON
            return JSONResponse(content={"status": "Processing failed", "error": str(e)}, status_code=500)